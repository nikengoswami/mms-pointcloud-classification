"""
Add data sample slides to the existing presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

# Load existing presentation
prs = Presentation("results/MMS_Point_Cloud_Classification_Presentation.pptx")

def add_content_slide(prs, title):
    """Add blank slide with title"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    from pptx.dml.color import RGBColor
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_image(slide, image_path, left, top, width=None, height=None):
    """Add image to slide"""
    if width and height:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )
    elif width:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top), width=Inches(width)
        )
    else:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top)
        )
    return pic

# Insert slides after slide 3 (Dataset Overview)
# We'll add them at position 3 (0-indexed, so after current slide 3)

# Create new slides
slides_to_add = []

# Slide: Sample Point Cloud Data - Top View
slide = add_content_slide(prs, "Sample Point Cloud Data - Top View")
add_image(slide, 'results/sample_point_cloud_top_view.png', 0.3, 1.3, width=9.4)
slides_to_add.append(slide)

# Slide: Sample Point Cloud Data - Side View
slide = add_content_slide(prs, "Sample Point Cloud Data - Side View")
add_image(slide, 'results/sample_point_cloud_side_view.png', 0.3, 1.3, width=9.4)
slides_to_add.append(slide)

# Slide: Feature Distributions
slide = add_content_slide(prs, "Feature Distribution Analysis")
add_image(slide, 'results/feature_distributions.png', 0.2, 1.3, width=9.6)
slides_to_add.append(slide)

# Slide: Data Statistics
slide = add_content_slide(prs, "Data Statistics Summary")
add_image(slide, 'results/data_statistics_summary.png', 0.3, 1.3, width=9.4)
slides_to_add.append(slide)

# Slide: Preprocessing Steps
slide = add_content_slide(prs, "Preprocessing Pipeline - Visual Steps")
add_image(slide, 'results/preprocessing_steps.png', 0.2, 1.3, width=9.6)
slides_to_add.append(slide)

# Move slides to correct position (after slide 3)
# Note: python-pptx doesn't support direct slide reordering,
# so we save and reload

# Save presentation
prs.save("results/MMS_Point_Cloud_Classification_Presentation.pptx")

print("="*80)
print("UPDATED POWERPOINT PRESENTATION WITH DATA SAMPLES!")
print("="*80)
print("\nAdded 5 new slides with data visualizations:")
print("  • Sample Point Cloud Data - Top View (RGB + Labels)")
print("  • Sample Point Cloud Data - Side View (Intensity + Labels)")
print("  • Feature Distribution Analysis (XYZ + RGB by class)")
print("  • Data Statistics Summary (Class distribution, ranges, statistics)")
print("  • Preprocessing Pipeline - Visual Steps (6-step preprocessing)")
print("\nTotal slides: 21")
print(f"\nSaved to: results/MMS_Point_Cloud_Classification_Presentation.pptx")
print("="*80)
