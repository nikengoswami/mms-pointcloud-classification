"""
Create PowerPoint presentation for MMS Point Cloud Classification project
Shows data samples, processing workflow, and results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_content_slide(prs, title):
    """Add blank slide with title"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False):
    """Add text box to slide"""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold

    return text_box

def add_bullet_points(slide, left, top, width, height, points, font_size=14):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(points):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = point
        p.level = 0
        p.font.size = Pt(font_size)

    return text_box

def add_image(slide, image_path, left, top, width=None, height=None):
    """Add image to slide"""
    if width and height:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )
    elif width:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top), width=Inches(width)
        )
    elif height:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top), height=Inches(height)
        )
    else:
        pic = slide.shapes.add_picture(
            str(image_path), Inches(left), Inches(top)
        )
    return pic


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
add_title_slide(
    prs,
    "MMS Point Cloud Classification",
    "Deep Learning for Semantic Segmentation\nDecember 2025"
)


# ============================================================================
# SLIDE 2: Project Overview
# ============================================================================
slide = add_content_slide(prs, "Project Overview")

add_text_box(slide, 0.5, 1.2, 4, 0.5, "Goal:", font_size=18, bold=True)
add_text_box(
    slide, 0.5, 1.6, 8.5, 0.8,
    "Develop AI-powered system to automatically classify Mobile Mapping System (MMS) "
    "point clouds into 5 semantic categories using deep learning.",
    font_size=14
)

add_text_box(slide, 0.5, 2.6, 4, 0.5, "Categories:", font_size=18, bold=True)
points = [
    "Road - Road surfaces, ground, bridge decks",
    "Snow - Snow coverage on surfaces",
    "Vehicle - Cars, trucks, and other vehicles",
    "Vegetation - Low, medium, and high vegetation",
    "Others - Buildings, unclassified objects, noise"
]
add_bullet_points(slide, 0.5, 3.0, 8.5, 2.5, points, font_size=14)

add_text_box(slide, 0.5, 5.7, 4, 0.5, "Status:", font_size=18, bold=True)
add_text_box(
    slide, 0.5, 6.1, 8.5, 0.5,
    "✅ COMPLETED - Achieved 94.78% accuracy (exceeds 88-90% target)",
    font_size=16, bold=True
)


# ============================================================================
# SLIDE 3: Dataset Statistics
# ============================================================================
slide = add_content_slide(prs, "Dataset Overview")

# Load data to get statistics
data_path = Path("data/processed/train_data.npz")
if data_path.exists():
    train_data = np.load(data_path)
    val_data = np.load("data/processed/val_data.npz")
    test_data = np.load("data/processed/test_data.npz")

    total_points = len(train_data['xyz']) + len(val_data['xyz']) + len(test_data['xyz'])

    stats_text = f"""Total Labeled Points: {total_points:,}

Data Split:
  • Training: {len(train_data['xyz']):,} points (70%)
  • Validation: {len(val_data['xyz']):,} points (15%)
  • Test: {len(test_data['xyz']):,} points (15%)

Features per Point: 7
  • XYZ Coordinates (3)
  • RGB Colors (3)
  • Intensity (1)

Source: CloudCompare-labeled LAS files"""

    add_text_box(slide, 0.5, 1.2, 4.5, 5.5, stats_text, font_size=16)

# Create class distribution chart
if data_path.exists():
    fig, ax = plt.subplots(figsize=(6, 4))

    # Count points per class
    all_labels = np.concatenate([
        train_data['labels'],
        val_data['labels'],
        test_data['labels']
    ])

    classes = ['Road', 'Snow', 'Vehicle', 'Vegetation', 'Others']
    class_counts = [np.sum(all_labels == i) for i in range(5)]

    colors = ['#E74C3C', '#3498DB', '#F39C12', '#27AE60', '#95A5A6']
    bars = ax.bar(classes, class_counts, color=colors, alpha=0.7, edgecolor='black', linewidth=1.5)

    ax.set_ylabel('Number of Points', fontsize=12, fontweight='bold')
    ax.set_title('Class Distribution', fontsize=14, fontweight='bold')
    ax.grid(axis='y', alpha=0.3)

    # Add value labels
    for bar, count in zip(bars, class_counts):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{count:,}',
                ha='center', va='bottom', fontsize=10, fontweight='bold')

    plt.xticks(rotation=15, ha='right')
    plt.tight_layout()
    plt.savefig('results/class_distribution.png', dpi=150, bbox_inches='tight')
    plt.close()

    add_image(slide, 'results/class_distribution.png', 5.2, 1.5, width=4.3)


# ============================================================================
# SLIDE 4: Data Processing Workflow - Part 1
# ============================================================================
slide = add_content_slide(prs, "Data Processing Workflow (1/2)")

workflow_text = """1. Data Collection & Labeling
   • Raw MMS point cloud data (LAS format)
   • Manual labeling in CloudCompare
   • Standard classification codes assigned

2. Data Loading & Parsing
   • Read LAS files using laspy library
   • Extract XYZ coordinates, RGB colors, intensity
   • Load classification labels from CloudCompare

3. Class Mapping
   • Map LAS classification codes to 5 target classes
   • Example mappings:
     - Code 2, 11, 17 → Road (Class 0)
     - Code 1 → Snow (Class 1)
     - Code 2 → Vehicle (Class 2)
     - Code 3, 4, 5 → Vegetation (Class 3)
     - Code 0, 6, 7, 9, 10 → Others (Class 4)"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, workflow_text, font_size=14)


# ============================================================================
# SLIDE 5: Data Processing Workflow - Part 2
# ============================================================================
slide = add_content_slide(prs, "Data Processing Workflow (2/2)")

workflow_text2 = """4. Data Splitting
   • Random shuffle of all labeled points
   • 70% Training, 15% Validation, 15% Test
   • Saved as compressed NumPy arrays (.npz)

5. Preprocessing for Training
   • Random sampling: 2048 points per batch
   • XYZ normalization: Zero-mean, unit variance
   • Feature concatenation: [X, Y, Z, R, G, B, Intensity]

6. Data Augmentation (Training only)
   • Random rotation around Z-axis (0-360°)
   • Random scaling (0.95-1.05x)
   • Improves model generalization

7. Batch Creation
   • Batch size: 8 for PointNet++, 16 for SimplePointNet
   • Parallel data loading with PyTorch DataLoader
   • GPU memory optimization"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, workflow_text2, font_size=14)


# ============================================================================
# SLIDE 6: Data Processing Flow Diagram
# ============================================================================
slide = add_content_slide(prs, "End-to-End Data Pipeline")

# Create flow diagram
fig, ax = plt.subplots(figsize=(10, 6))
ax.set_xlim(0, 10)
ax.set_ylim(0, 8)
ax.axis('off')

# Box properties
box_width = 2.2
box_height = 0.8

def draw_box(x, y, text, color='lightblue'):
    """Draw a box with text"""
    rect = mpatches.FancyBboxPatch(
        (x - box_width/2, y - box_height/2),
        box_width, box_height,
        boxstyle="round,pad=0.1",
        edgecolor='black',
        facecolor=color,
        linewidth=2
    )
    ax.add_patch(rect)
    ax.text(x, y, text, ha='center', va='center',
            fontsize=10, fontweight='bold', wrap=True)

def draw_arrow(x1, y1, x2, y2):
    """Draw arrow between boxes"""
    ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle='->', lw=2, color='black'))

# Draw boxes
draw_box(1.5, 7, 'Raw LAS File\n(MMS Data)', '#FFE5E5')
draw_box(1.5, 5.5, 'CloudCompare\nLabeling', '#E5F5FF')
draw_box(1.5, 4, 'Class Mapping\n& Filtering', '#E5FFE5')
draw_box(1.5, 2.5, 'Train/Val/Test\nSplit', '#FFF5E5')

draw_box(5, 7, 'Preprocessed\nNumPy Arrays', '#FFE5E5')
draw_box(5, 5.5, 'Random\nSampling', '#E5F5FF')
draw_box(5, 4, 'Normalization\n& Augmentation', '#E5FFE5')
draw_box(5, 2.5, 'Batch\nCreation', '#FFF5E5')

draw_box(8.5, 5.5, 'PointNet++\nModel', '#E5E5FF')
draw_box(8.5, 4, 'Training\n(30 epochs)', '#FFE5FF')
draw_box(8.5, 2.5, 'Evaluation\n& Metrics', '#E5FFE5')

# Draw arrows - vertical
draw_arrow(1.5, 6.6, 1.5, 5.9)
draw_arrow(1.5, 5.1, 1.5, 4.4)
draw_arrow(1.5, 3.6, 1.5, 2.9)

draw_arrow(5, 6.6, 5, 5.9)
draw_arrow(5, 5.1, 5, 4.4)
draw_arrow(5, 3.6, 5, 2.9)

draw_arrow(8.5, 5.1, 8.5, 4.4)
draw_arrow(8.5, 3.6, 8.5, 2.9)

# Draw arrows - horizontal
draw_arrow(2.6, 2.5, 3.9, 2.5)
draw_arrow(2.6, 7, 3.9, 7)
draw_arrow(6.1, 5.5, 7.4, 5.5)
draw_arrow(6.1, 4, 7.4, 4)

# Add labels
ax.text(3.25, 7.3, 'Process', ha='center', fontsize=8, style='italic')
ax.text(3.25, 2.8, 'Save', ha='center', fontsize=8, style='italic')
ax.text(6.75, 5.8, 'Train', ha='center', fontsize=8, style='italic')
ax.text(6.75, 4.3, 'Optimize', ha='center', fontsize=8, style='italic')

plt.tight_layout()
plt.savefig('results/data_flow_diagram.png', dpi=150, bbox_inches='tight', facecolor='white')
plt.close()

add_image(slide, 'results/data_flow_diagram.png', 0.3, 1.5, width=9.4)


# ============================================================================
# SLIDE 7: Model Architecture Comparison
# ============================================================================
slide = add_content_slide(prs, "Model Architecture Comparison")

arch_text = """SimplePointNet (Baseline)
• Single-scale global feature extraction
• Parameters: 192,517
• Architecture: MLP + MaxPooling
• Training Time: ~2.5 hours (GPU)
• Accuracy: 86.01%

PointNet++ (Final Model)
• Multi-scale hierarchical feature learning
• Parameters: 968,069
• Architecture: 4 Set Abstraction + 4 Feature Propagation layers
• Training Time: ~4 hours (GPU)
• Accuracy: 94.78% ✅

Key Difference:
PointNet++ captures both local and global features at multiple scales,
enabling better understanding of complex 3D structures."""

add_text_box(slide, 0.5, 1.2, 9, 5.8, arch_text, font_size=14)


# ============================================================================
# SLIDE 8: Training Configuration
# ============================================================================
slide = add_content_slide(prs, "Training Configuration")

config_text = """Model: PointNet++
• Batch Size: 8
• Points per Sample: 2048
• Learning Rate: 0.001 (Adam optimizer)
• Epochs: 30
• Scheduler: ReduceLROnPlateau (patience=5, factor=0.5)

Data Augmentation:
• Random rotation around Z-axis (0-360°)
• Random scaling (0.95-1.05x)

Hardware:
• GPU: NVIDIA GeForce RTX 4050 Laptop GPU
• PyTorch: 2.5.1+cu121 (CUDA-enabled)
• GPU Utilization: ~100% during training
• Training Duration: ~4 hours for 30 epochs

Loss Function:
• CrossEntropyLoss for multi-class classification"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, config_text, font_size=14)


# ============================================================================
# SLIDE 9: Final Results - Overall Metrics
# ============================================================================
slide = add_content_slide(prs, "Final Results: Overall Performance")

# Add overall metrics comparison image
if Path('results/overall_metrics_comparison.png').exists():
    add_image(slide, 'results/overall_metrics_comparison.png', 0.5, 1.5, width=5)

# Add metrics table
metrics_text = """PointNet++ Test Results:
• Overall Accuracy: 94.78%
• Mean IoU: 87.51%
• Kappa Coefficient: 0.9187
• Weighted F1-Score: 0.9479

Improvement over SimplePointNet:
• Accuracy: +8.77%
• Mean IoU: +11.72%
• Kappa: +0.1445"""

add_text_box(slide, 6, 1.5, 3.5, 4, metrics_text, font_size=14, bold=False)


# ============================================================================
# SLIDE 10: Per-Class Performance
# ============================================================================
slide = add_content_slide(prs, "Final Results: Per-Class Performance")

# Add per-class IoU comparison
if Path('results/per_class_iou_comparison.png').exists():
    add_image(slide, 'results/per_class_iou_comparison.png', 0.3, 1.3, width=9.4)


# ============================================================================
# SLIDE 11: Confusion Matrix
# ============================================================================
slide = add_content_slide(prs, "Final Results: Confusion Matrix")

# Add confusion matrix
if Path('results/pointnet2_confusion_matrix_normalized.png').exists():
    add_image(slide, 'results/pointnet2_confusion_matrix_normalized.png', 1, 1.5, width=8)


# ============================================================================
# SLIDE 12: Key Achievements
# ============================================================================
slide = add_content_slide(prs, "Key Achievements")

achievements_text = """Technical Achievements:
✅ Fixed PointNet++ dimension mismatch bugs
✅ Implemented GPU acceleration (6x speedup)
✅ Achieved 94.78% accuracy (exceeds 88-90% target)
✅ Comprehensive evaluation metrics (IoU, Kappa, F1)
✅ Complete data pipeline from LAS to predictions

Performance Achievements:
✅ Excellent overall accuracy: 94.78%
✅ Strong mean IoU: 87.51%
✅ Excellent Kappa coefficient: 0.9187
✅ All classes achieve >79% IoU

Best Performing Classes:
• Snow: 91.87% IoU (+20.37% vs SimplePointNet)
• Road: 91.45% IoU
• Others: 89.75% IoU
• Vegetation: 85.30% IoU (+24.04% improvement!)"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, achievements_text, font_size=14)


# ============================================================================
# SLIDE 13: Challenges & Solutions
# ============================================================================
slide = add_content_slide(prs, "Challenges & Solutions")

challenges_text = """Challenge 1: PointNet++ Dimension Mismatch
Problem: Expected 7 channels but got 10 channels
Solution: Changed in_channel from num_features (7) to num_features + 3 (10)
Status: ✅ Fixed

Challenge 2: Tensor Format Mismatch
Problem: Index out of bounds in query_ball_point function
Solution: Added tensor permutations between encoder/decoder layers
Status: ✅ Fixed

Challenge 3: GPU Support
Problem: PyTorch CPU-only installation, no CUDA acceleration
Solution: Installed PyTorch 2.5.1+cu121 with CUDA support
Impact: Training time reduced from 24+ hours to 4 hours (6x speedup)
Status: ✅ Fixed

Challenge 4: RandLA-Net Implementation
Problem: torch.gather() dimension mismatch
Status: ⚠️ Identified but postponed (not critical for project goals)"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, challenges_text, font_size=13)


# ============================================================================
# SLIDE 14: Improvement Heatmap
# ============================================================================
slide = add_content_slide(prs, "Model Improvement Analysis")

# Add improvement heatmap
if Path('results/improvement_heatmap.png').exists():
    add_image(slide, 'results/improvement_heatmap.png', 0.5, 1.5, width=9)


# ============================================================================
# SLIDE 15: Conclusion & Future Work
# ============================================================================
slide = add_content_slide(prs, "Conclusion & Future Work")

conclusion_text = """Project Status: ✅ SUCCESSFULLY COMPLETED

Summary:
• Developed AI-powered MMS point cloud classification system
• Achieved 94.78% accuracy (exceeds 88-90% target by 4.78-6.78%)
• Mean IoU: 87.51%, Kappa: 0.9187 (excellent agreement)
• Complete pipeline from raw LAS files to predictions
• Comprehensive evaluation and documentation

Recommended Model: PointNet++ (checkpoints/pointnet2_best_model.pth)

Future Improvements (Optional):
• Fix RandLA-Net implementation for comparison
• Enhanced data augmentation (point dropout, Gaussian noise)
• Weighted loss for class imbalance (Vehicle: only 4,836 points)
• Model optimization (quantization, TorchScript, ONNX export)
• Real-time inference pipeline for large point clouds
• Ensemble methods combining multiple models"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, conclusion_text, font_size=14)


# ============================================================================
# SLIDE 16: References & Contact
# ============================================================================
slide = add_content_slide(prs, "References & Documentation")

references_text = """Key References:
1. PointNet++: Qi et al. (2017), NeurIPS 2017
2. PointNet: Qi et al. (2017), CVPR 2017
3. RandLA-Net: Hu et al. (2020), CVPR 2020
4. ASPRS LAS 1.4 Format Specification
5. CloudCompare: https://www.cloudcompare.org/

Project Documentation:
• README.md - Complete project overview
• FINAL_PROJECT_SUMMARY.md - Detailed technical summary
• results/model_comparison.md - Model comparison analysis
• All code, models, and results available in repository

Dataset:
• Total: 1,461,189 labeled points
• Source: CloudCompare-labeled MMS LAS files
• 5 semantic categories: Road, Snow, Vehicle, Vegetation, Others

Timeline:
• Started: December 2025
• Completed: December 25, 2025
• Total Duration: ~1 week"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, references_text, font_size=14)


# Save presentation
output_path = Path("results/MMS_Point_Cloud_Classification_Presentation.pptx")
prs.save(str(output_path))

print("="*80)
print("POWERPOINT PRESENTATION CREATED SUCCESSFULLY!")
print("="*80)
print(f"\nSaved to: {output_path}")
print(f"\nPresentation contains 16 slides:")
print("  1. Title Slide")
print("  2. Project Overview")
print("  3. Dataset Statistics")
print("  4. Data Processing Workflow (Part 1)")
print("  5. Data Processing Workflow (Part 2)")
print("  6. End-to-End Data Pipeline Diagram")
print("  7. Model Architecture Comparison")
print("  8. Training Configuration")
print("  9. Final Results: Overall Performance")
print(" 10. Final Results: Per-Class Performance")
print(" 11. Final Results: Confusion Matrix")
print(" 12. Key Achievements")
print(" 13. Challenges & Solutions")
print(" 14. Model Improvement Analysis")
print(" 15. Conclusion & Future Work")
print(" 16. References & Documentation")
print("\n" + "="*80)
